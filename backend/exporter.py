import os
import uuid
from typing import List, Dict, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)

# Basic palette
PRIMARY = RGBColor(16, 185, 129)   # emerald
ACCENT = RGBColor(245, 158, 11)     # amber
DARK = RGBColor(12, 17, 23)
LIGHT = RGBColor(240, 246, 252)


def _set_font(paragraph, size=24, bold=False, color=LIGHT, align=PP_ALIGN.LEFT):
    paragraph.alignment = align
    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_textbox(slide, text, left, top, width, height, size=20, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text or ""
    _set_font(p, size=size, bold=bold, align=align)
    return box


def _add_title_block(slide, title: str, subtitle: Optional[str], presenter: Optional[str]):
    _add_textbox(slide, title or "Title", Inches(1), Inches(1), SLIDE_WIDTH - Inches(2), Inches(2), size=44, bold=True, align=PP_ALIGN.CENTER)
    if subtitle:
        _add_textbox(slide, subtitle, Inches(2), Inches(2.6), SLIDE_WIDTH - Inches(4), Inches(1.2), size=24, align=PP_ALIGN.CENTER)
    if presenter:
        _add_textbox(slide, presenter, Inches(3.5), Inches(3.4), SLIDE_WIDTH - Inches(7), Inches(0.8), size=18, align=PP_ALIGN.CENTER)


def _add_bullets(slide, title: str, bullets: List[str]):
    _add_textbox(slide, title, Inches(1), Inches(0.8), SLIDE_WIDTH - Inches(2), Inches(1), size=32, bold=True)
    top = Inches(1.8)
    for b in bullets:
        box = slide.shapes.add_textbox(Inches(1.2), top, SLIDE_WIDTH - Inches(2.4), Inches(0.7))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = f"• {b}" if b else "•"
        _set_font(p, size=20)
        top += Inches(0.8)


def create_pptx(slides: List[Dict], title: str = "Presentation") -> str:
    prs = Presentation()
    prs.slide_height = SLIDE_HEIGHT
    prs.slide_width = SLIDE_WIDTH
    blank = prs.slide_layouts[6]

    for slide_data in slides:
        layout = slide_data.get("layout", "title_slide")
        s = prs.slides.add_slide(blank)

        if layout in ["title_slide", "title"]:
            _add_title_block(s, slide_data.get("title"), slide_data.get("tagline"), slide_data.get("presenter"))
        elif layout == "overview":
            bullets = [slide_data.get("definition"), slide_data.get("point1"), slide_data.get("point2"), slide_data.get("point3")]
            _add_bullets(s, slide_data.get("title", "Overview"), [b for b in bullets if b])
        elif layout == "problem":
            _add_bullets(s, slide_data.get("title", "Problem Statement"), [slide_data.get("problem"), slide_data.get("challenge1"), slide_data.get("challenge2"), slide_data.get("why_matters")])
        elif layout == "background":
            _add_bullets(s, slide_data.get("title", "Background & Context"), [slide_data.get("description"), slide_data.get("milestone1"), slide_data.get("milestone2"), slide_data.get("current_state")])
        elif layout == "concepts":
            _add_bullets(s, slide_data.get("title", "Core Concepts"), [slide_data.get("concept1"), slide_data.get("concept2"), slide_data.get("concept3"), slide_data.get("concept4")])
        elif layout == "process":
            _add_bullets(s, slide_data.get("title", "How It Works"), [slide_data.get("step1"), slide_data.get("step2"), slide_data.get("step3"), slide_data.get("explanation")])
        elif layout == "applications":
            _add_bullets(s, slide_data.get("title", "Applications & Use Cases"), [slide_data.get("use_case1"), slide_data.get("use_case2"), slide_data.get("use_case3"), slide_data.get("use_case4")])
        elif layout == "benefits":
            _add_bullets(s, slide_data.get("title", "Benefits"), [slide_data.get("benefit1"), slide_data.get("benefit2"), slide_data.get("benefit3"), slide_data.get("value")])
        elif layout == "challenges":
            _add_bullets(s, slide_data.get("title", "Challenges & Limitations"), [slide_data.get("limitation1"), slide_data.get("limitation2"), slide_data.get("risk"), slide_data.get("improvement")])
        elif layout == "conclusion":
            _add_bullets(s, slide_data.get("title", "Conclusion"), [slide_data.get("summary"), slide_data.get("future"), slide_data.get("closing")])
        else:
            # Fallback
            _add_bullets(s, slide_data.get("title", "Slide"), [slide_data.get("content")])

    os.makedirs("exports", exist_ok=True)
    file_id = uuid.uuid4().hex
    out_path = os.path.join("exports", f"{file_id}.pptx")
    prs.save(out_path)
    return out_path
